import streamlit as st
import pandas as pd
import numpy as np
import re
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

# =========================================================================
# ⚙️ INITIALIZATION & CONFIGURATION
# =========================================================================
st.set_page_config(
    page_title="Infrastructure Asset Dashboard",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

SPREADSHEET_ID = "1FGKOzWoUrbf3PXN_ahgG1t-83JZT4H4sioQepePbBxM"
csv_url = f"https://docs.google.com/spreadsheets/d/{SPREADSHEET_ID}/gviz/tq?tqx=out:csv"

# 🔥 ENGINE PERBAIKAN LINK GOOGLE DRIVE (MENDUKUNG SEMUA FORMAT LINK DRIVE)
def dapatkan_link_gambar_drive(input_url):
    val = str(input_url).strip()
    if val == "" or val.lower() == "nan" or "http" not in val.lower():
        # Jika teks di sheet ternyata sudah berupa ID murni sepanjang 33 karakter tanpa embel-embel
        if len(val) >= 25 and "/" not in val:
            return f"https://drive.google.com/uc?export=view&id={val}"
        return None
    
    # regex untuk menangkap ID unik file Google Drive (Kombinasi huruf besar, kecil, angka, strip, underscore)
    id_match = re.search(r'(?:id=|/d/|/file/d/)([a-zA-Z0-9_-]{25,45})', val)
    if id_match:
        file_id = id_match.group(1)
        return f"https://drive.google.com/uc?export=view&id={file_id}"
    
    # Fallback terakhir jika regex gagal namun link berisi pola ID standar
    match_fallback = re.search(r'([a-zA-Z0-9_-]{33})', val)
    if match_fallback:
        return f"https://drive.google.com/uc?export=view&id={match_fallback.group(1)}"
        
    return val

# Custom CSS Modern untuk membuat struktur layout huruf, angka, dan frame foto terlihat sangat rapi
st.markdown("""
    <style>
    .report-card {
        background-color: #f8fafc;
        border: 1px solid #e2e8f0;
        border-radius: 12px;
        padding: 24px;
        margin-bottom: 25px;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.05);
    }
    .site-title {
        color: #0f172a;
        font-size: 22px;
        font-weight: 700;
        margin-bottom: 15px;
        border-bottom: 2px solid #10b981;
        padding-bottom: 5px;
    }
    .info-label {
        font-weight: 600;
        color: #475569;
        display: inline-block;
        width: 150px;
    }
    .info-value {
        color: #0f172a;
    }
    .image-container {
        border: 2px solid #cbd5e1;
        border-radius: 8px;
        overflow: hidden;
        background-color: #ffffff;
        text-align: center;
        padding: 5px;
    }
    </style>
""", unsafe_allow_html=True)

st.title("⚡ Telecommunication Infrastructure Dashboard")
st.caption("Sinkronisasi Otomatis Google Sheets Resource (348 Tsel) & Rendering Gambar Google Drive Terpilih")
st.markdown("---")

# =========================================================================
# 💾 LOAD DATA RESOURCE
# =========================================================================
@st.cache_data(ttl=2)
def load_data_sheets(url):
    try:
        df = pd.read_csv(url)
        return df
    except Exception as e:
        return str(e)

with st.spinner("🔄 Sedang mengunduh dan merapikan struktur data tabel..."):
    df_raw = load_data_sheets(csv_url)

if isinstance(df_raw, str):
    st.error(f"❌ Koneksi ke Google Sheets Gagal. Detail Error: `{df_raw}`")
    st.stop()
elif df_raw is None or df_raw.empty:
    st.warning("⚠️ Google Sheets terhubung namun tidak mendeteksi adanya baris data.")
    st.stop()
else:
    # Lakukan standarisasi pembersihan spasi di nama kolom
    df_raw.columns = [str(c).strip() for c in df_raw.columns]
    df = df_raw.dropna(how='all', axis=0).copy()

    # =========================================================================
    # 🔍 DETEKSI NAMA KOLOM SPREADSHEET (MENDUKUNG SEGALA VARIASI HURUF)
    # =========================================================================
    KOLOM_UTAMA = next((c for c in df.columns if any(p in c.upper() for p in ['SITE ID', 'SITE_ID', 'SITEID', 'TOWER ID', 'COMPONENT'])), df.columns[0])
    KOLOM_FOTO = next((c for c in df.columns if any(p in c.upper() for p in ['LINK FOTO', 'DOKUMENTASI', 'DRIVE', 'FOTO', 'URL FOTO', 'GAMBAR', 'LINK'])), None)
    KOLOM_VR = next((c for c in df.columns if any(p in c.upper() for p in ['VOLTAGE R', 'TEGANGAN R', 'VR', 'V_R'])), None)
    KOLOM_VS = next((c for c in df.columns if any(p in c.upper() for p in ['VOLTAGE S', 'TEGANGAN S', 'VS', 'V_S'])), None)
    KOLOM_VT = next((c for c in df.columns if any(p in c.upper() for p in ['VOLTAGE T', 'TEGANGAN T', 'VT', 'V_T'])), None)

    # --- PANEL PENCARIAN DI SIDEBAR ---
    st.sidebar.header("🔍 Sistem Pencarian Asset")
    df[KOLOM_UTAMA] = df[KOLOM_UTAMA].fillna("KOSONG").astype(str).str.strip()
    
    search_input = st.sidebar.text_input("Ketik Site ID di Sini (Bisa Ketik Sebagian):", "").strip()
    
    if search_input:
        df_filtered = df[df[KOLOM_UTAMA].str.contains(search_input, case=False, na=False)]
    else:
        list_pilihan = sorted(list(df[KOLOM_UTAMA].unique()))
        selected_id = st.sidebar.selectbox("Atau Pilih ID Sesuai List Dropdown:", list_pilihan)
        df_filtered = df[df[KOLOM_UTAMA] == selected_id]

    if st.sidebar.button("🔄 Bersihkan Sisa Cache & Refresh", use_container_width=True):
        st.cache_data.clear()
        st.rerun()

    # =========================================================================
    # 🏗️ LAYOUT 3 BAGIAN MODEL KARTU (SANGAT RAPI & PRESISI)
    # =========================================================================
    if df_filtered.empty:
        st.warning("⚠️ Maaf, Site ID atau data asset yang Anda cari tidak ditemukan.")
    else:
        for idx, row in df_filtered.reset_index(drop=True).iterrows():
            
            # Membuka pembungkus box kartu HTML
            st.markdown(f'<div class="report-card">', unsafe_allow_html=True)
            st.markdown(f'<div class="site-title">📍 SITE ID: {row[KOLOM_UTAMA]}</div>', unsafe_allow_html=True)
            
            # Membuat layout 3 kolom sejajar dengan rasio proporsional lebar layar
            col_teks, col_grafik, col_foto = st.columns([1.6, 1.2, 1.2])
            
            # --- BAGIAN 1: DETAIL KETERANGAN HURUF & ANGKA ---
            with col_teks:
                st.markdown("📋 **Informasi Spesifikasi Teknis:**")
                html_konten_teks = '<div style="line-height: 1.8; font-size:14px;">'
                
                for kolom_nama in df_filtered.columns:
                    if kolom_nama not in [KOLOM_UTAMA, KOLOM_FOTO, KOLOM_VR, KOLOM_VS, KOLOM_VT]:
                        isi_sel = str(row[kolom_nama]).strip()
                        if isi_sel != "" and isi_sel.lower() != "nan":
                            html_konten_teks += f'<div><span class="info-label">• {kolom_nama}</span>: <span class="info-value">{isi_sel}</span></div>'
                            
                html_konten_teks += '</div>'
                st.markdown(html_konten_teks, unsafe_allow_html=True)

            # --- BAGIAN 2: POLA ARUS POWER VOLTAGE RST ---
            with col_grafik:
                st.markdown("📊 **Analisa Pola Arus Tegangan:**")
                
                # Mengambil nilai fasa voltase
                v_r = row[KOLOM_VR] if KOLOM_VR and pd.notnull(row[KOLOM_VR]) else "0"
                v_s = row[KOLOM_VS] if KOLOM_VS and pd.notnull(row[KOLOM_VS]) else "0"
                v_t = row[KOLOM_VT] if KOLOM_VT and pd.notnull(row[KOLOM_VT]) else "0"
                
                # Menampilkan nilai angka fasa dalam bentuk sub-metrik kecil
                m_c1, m_c2, m_c3 = st.columns(3)
                m_c1.metric("R-Phase", f"{v_r} V")
                m_c2.metric("S-Phase", f"{v_s} V")
                m_c3.metric("T-Phase", f"{v_t} V")
                
                # Membuat Grafik Pola Kestabilan Daya Batang Mini
                try:
                    num_r = float(str(v_r).replace('V','').strip() or 0)
                    num_s = float(str(v_s).replace('V','').strip() or 0)
                    num_t = float(str(v_t).replace('V','').strip() or 0)
                    data_grafik_volt = pd.DataFrame({"Voltase (V)": [num_r, num_s, num_t]}, index=["Fasa R", "Fasa S", "Fasa T"])
                    st.bar_chart(data_grafik_volt, height=130)
                except:
                    pass

            # --- BAGIAN 3: LAYOUT KHUSUS GAMBAR GOOGLE DRIVE (WAJIB TAMPIL) ---
            with col_foto:
                st.markdown("🖼️ **Dokumentasi Visual Lapangan:**")
                link_mentah_foto = str(row[KOLOM_FOTO]).strip() if KOLOM_FOTO else ""
                
                if link_mentah_foto != "" and link_mentah_foto.lower() != "nan":
                    url_gambar_matang = dapatkan_link_gambar_drive(link_mentah_foto)
                    
                    if url_gambar_matang:
                        # Tampilan frame khusus foto menggunakan standar HTML responsive bootstrap style
                        st.markdown(f"""
                            <div class="image-container">
                                <img src="{url_gambar_matang}" style="width:100%; max-height:220px; object-fit:contain; border-radius:4px;" 
                                     onerror="this.onerror=null; this.src='https://placehold.co/400x300?text=Akses+Drive+Dibatasi+Sila+Cek+Privacy';">
                            </div>
                        """, unsafe_allow_html=True)
                        st.caption(f"<center>Foto Validasi - {row[KOLOM_UTAMA]}</center>", unsafe_allow_html=True)
                else:
                    st.error("❌ Kolom link foto kosong / belum diinput di spreadsheet.")
            
            # Menutup box kartu HTML
            st.markdown('</div>', unsafe_allow_html=True)

        # =========================================================================
        # 📥 ENGINE GENERATOR PPTX MERANGKUM LENGKAP DETAIL DAN GAMBAR DRIVE
        # =========================================================================
        st.markdown("<br><br>", unsafe_allow_html=True)
        st.subheader("📥 Penarikan Dokumen Laporan Akhir (.PPTX)")
        st.info("💡 *Tombol di bawah ini akan mengekspor seluruh data yang tampil di atas ke dalam slide PowerPoint lengkap beserta fotonya.*")
        
        def buat_dokumen_pptx_lengkap(data_sumber):
            prs = Presentation()
            
            for _, baris_data in data_sumber.iterrows():
                # Menambahkan slide kosong (Blank Layout) per baris Site ID
                slide_obj = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 1. Judul Atas Slide
                box_judul = slide_obj.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                tf_judul = box_judul.text_frame
                para_judul = tf_judul.paragraphs[0]
                para_judul.text = f"LAPORAN DATA INFRASTRUKTUR: {baris_data[KOLOM_UTAMA]}"
                para_judul.font.size = Pt(20)
                para_judul.font.bold = True
                
                # 2. Penyusunan Keterangan Huruf & Angka (Sisi Kiri Slide PPT)
                box_teks = slide_obj.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.3))
                tf_teks = box_teks.text_frame
                tf_teks.word_wrap = True
                
                # Baris Informasi Tegangan Voltase di posisi paling atas teks slide
                para_v = tf_teks.paragraphs[0]
                v_r_slide = baris_data[KOLOM_VR] if KOLOM_VR else "-"
                v_s_slide = baris_data[KOLOM_VS] if KOLOM_VS else "-"
                v_t_slide = baris_data[KOLOM_VT] if KOLOM_VT else "-"
                para_v.text = f"Pola Tegangan RST: R={v_r_slide}V | S={v_s_slide}V | T={v_t_slide}V\n"
                para_v.font.bold = True
                para_v.font.size = Pt(12)
                
                # Perulangan memasukkan sisa kolom huruf dan angka pendukung
                for nama_kol_data in data_sumber.columns:
                    if nama_kol_data not in [KOLOM_UTAMA, KOLOM_FOTO, KOLOM_VR, KOLOM_VS, KOLOM_VT]:
                        teks_isi_kolom = str(baris_data[nama_kol_data]).strip()
                        if teks_isi_kolom != "" and teks_isi_kolom.lower() != "nan":
                            para_info = tf_teks.add_paragraph()
                            para_info.text = f"• {nama_kol_data}: {teks_isi_kolom}"
                            para_info.font.size = Pt(11)
                
                # 3. Proses Penarikan Otomatis Gambar dari Drive untuk Disisipkan ke Slide Sisi Kanan
                cell_alamat_foto = str(baris_data[KOLOM_FOTO]).strip() if KOLOM_FOTO else ""
                if cell_alamat_foto != "" and cell_alamat_foto.lower() != "nan":
                    tautan_langsung_foto = dapatkan_link_gambar_drive(cell_alamat_foto)
                    if tautan_langsung_foto:
                        try:
                            # Server melakukan request download byte file gambar ke Google Drive
                            koneksi_request = requests.get(tautan_langsung_foto, timeout=8)
                            if koneksi_request.status_code == 200:
                                bytes_gambar = BytesIO(koneksi_request.content)
                                # Memasukkan foto ke sisi kanan slide (Left=5.2 inci, Top=1.3 inci, Width=4.3 inci)
                                slide_obj.shapes.add_picture(bytes_gambar, Inches(5.2), Inches(1.3), width=Inches(4.3))
                        except Exception:
                            # Jika link bermasalah atau terkena enkripsi Restricted, buat penanda teks di slide
                            box_err = slide_obj.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(2))
                            box_err.text_frame.text = "[Gambar tidak dapat disalin ke PPT karena hak akses file di Google Drive Anda diatur ke 'Dibatasi / Restricted']"
            
            output_stream = BytesIO()
            prs.save(output_stream)
            output_stream.seek(0)
            return output_stream

        # Tombol download untuk menarik file PPTX laporan lengkap
        if st.button("📊 Ambil File PowerPoint Laporan Lengkap (.PPTX)", type="primary", use_container_width=True):
            with st.spinner("⏳ Sedang menyusun struktur teks dan menyalin foto dari Google Drive ke slide PowerPoint Anda..."):
                file_final_ppt = buat_dokumen_pptx_lengkap(df_filtered)
                st.download_button(
                    label="📥 Klik di Sini Untuk Mengunduh File PPTX Anda",
                    data=file_final_ppt,
                    file_name=f"Laporan_Komprehensif_Site_{search_input or 'Selected'}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
